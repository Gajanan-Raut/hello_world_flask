# from flask import Flask, render_template

# app = Flask(__name__)

# @app.route("/")
# def home():
#     return render_template("index.html")

# if __name__ == "__main__":
#     app.run(debug=True)

from flask import Flask, render_template, request, send_file, jsonify
import os
from werkzeug.utils import secure_filename
from io import BytesIO

# PDF libraries
import PyPDF2
import fitz  # PyMuPDF
import pikepdf
from pdf2docx import Converter
from PIL import Image
from pdf2image import convert_from_path
import tempfile

# Optional/advanced libraries (import when installed)
try:
    import camelot
except Exception:
    camelot = None

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

ALLOWED = set(['pdf', 'png', 'jpg', 'jpeg'])

def allowed(filename):
    return '.' in filename and filename.rsplit('.',1)[1].lower() in ALLOWED

@app.route('/')
def index():
    return render_template('index.html')

# Generic upload endpoint used by the UI
@app.route('/upload', methods=['POST'])
def upload():
    f = request.files.get('file')
    if not f or not allowed(f.filename):
        return jsonify({'error':'No file or invalid extension'}), 400
    fname = secure_filename(f.filename)
    path = os.path.join(app.config['UPLOAD_FOLDER'], fname)
    f.save(path)
    return jsonify({'filename': fname})

# Merge PDFs: accept multiple files
@app.route('/merge', methods=['POST'])
def merge_pdfs():
    files = request.files.getlist('files')
    merger = PyPDF2.PdfMerger()
    for f in files:
        merger.append(f)
    out = BytesIO()
    merger.write(out)
    out.seek(0)
    return send_file(out, as_attachment=True, download_name='merged.pdf', mimetype='application/pdf')

# Split PDF: split into single-page PDFs and return zip (for simplicity, return first page as demo)
@app.route('/split', methods=['POST'])
def split_pdf():
    f = request.files.get('file')
    reader = PyPDF2.PdfReader(f)
    outputs = []
    for i, pg in enumerate(reader.pages):
        writer = PyPDF2.PdfWriter()
        writer.add_page(pg)
        buf = BytesIO()
        writer.write(buf)
        buf.seek(0)
        outputs.append((f'page_{i+1}.pdf', buf.read()))
    # For demo, return first page
    return send_file(BytesIO(outputs[0][1]), as_attachment=True, download_name=outputs[0][0], mimetype='application/pdf')

# Compress PDF using pikepdf (linearize + recompress images)
@app.route('/compress', methods=['POST'])
def compress_pdf():
    f = request.files.get('file')
    in_mem = BytesIO(f.read())
    in_mem.seek(0)
    out_mem = BytesIO()
    with pikepdf.Pdf.open(in_mem) as pdf:
        pdf.save(out_mem, optimize_streams=True, linearize=True)
    out_mem.seek(0)
    return send_file(out_mem, as_attachment=True, download_name='compressed.pdf', mimetype='application/pdf')

# PDF -> DOCX (uses pdf2docx)
@app.route('/pdf2docx', methods=['POST'])
def pdf_to_docx():
    f = request.files.get('file')
    src = os.path.join(app.config['UPLOAD_FOLDER'], secure_filename(f.filename))
    f.save(src)
    target = src + '.docx'
    cv = Converter(src)
    cv.convert(target, start=0, end=None)
    cv.close()
    return send_file(target, as_attachment=True)

# PDF -> PPTX (convert pages to images then build pptx)
@app.route('/pdf2pptx', methods=['POST'])
def pdf_to_pptx():
    from pptx import Presentation
    from pptx.util import Inches
    f = request.files.get('file')
    src = os.path.join(app.config['UPLOAD_FOLDER'], secure_filename(f.filename))
    f.save(src)
    images = convert_from_path(src, dpi=150)
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    for img in images:
        slide = prs.slides.add_slide(blank_slide_layout)
        img_io = BytesIO()
        img.save(img_io, format='PNG')
        img_io.seek(0)
        slide.shapes.add_picture(img_io, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
    out = BytesIO()
    prs.save(out)
    out.seek(0)
    return send_file(out, as_attachment=True, download_name='presentation.pptx')

# PDF -> JPG: return first page as JPG
@app.route('/pdf2jpg', methods=['POST'])
def pdf_to_jpg():
    f = request.files.get('file')
    images = convert_from_path(f.filename if False else BytesIO(f.read()), dpi=150)
    # convert_from_path expects path; as a simple approach, save temporarily
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
    tmp.write(f.read())
    tmp.flush()
    imgs = convert_from_path(tmp.name, dpi=150)
    img_io = BytesIO()
    imgs[0].save(img_io, format='JPEG')
    img_io.seek(0)
    os.unlink(tmp.name)
    return send_file(img_io, as_attachment=True, download_name='page1.jpg', mimetype='image/jpeg')

# JPG -> PDF
@app.route('/jpg2pdf', methods=['POST'])
def jpg_to_pdf():
    f = request.files.get('file')
    img = Image.open(f.stream).convert('RGB')
    out = BytesIO()
    img.save(out, format='PDF')
    out.seek(0)
    return send_file(out, as_attachment=True, download_name='converted.pdf', mimetype='application/pdf')

# Add watermark (text) — simple overlay using PyMuPDF
@app.route('/watermark', methods=['POST'])
def watermark_pdf():
    f = request.files.get('file')
    text = request.form.get('text', 'Watermark')
    doc = fitz.open(stream=f.read(), filetype='pdf')
    for page in doc:
        rect = page.rect
        page.insert_text((rect.width/4, rect.height/2), text, fontsize=40, rotate=45, color=(0.7,0.7,0.7), render_mode=3)
    out = BytesIO()
    doc.save(out)
    out.seek(0)
    return send_file(out, as_attachment=True, download_name='watermarked.pdf', mimetype='application/pdf')

# Rotate PDF
@app.route('/rotate', methods=['POST'])
def rotate_pdf():
    degrees = int(request.form.get('degrees', '90'))
    f = request.files.get('file')
    reader = PyPDF2.PdfReader(f)
    writer = PyPDF2.PdfWriter()
    for p in reader.pages:
        p.rotate_clockwise(degrees)
        writer.add_page(p)
    out = BytesIO()
    writer.write(out)
    out.seek(0)
    return send_file(out, as_attachment=True, download_name='rotated.pdf', mimetype='application/pdf')

# Protect PDF (add owner/user password) using PyPDF2
@app.route('/protect', methods=['POST'])
def protect_pdf():
    pwd = request.form.get('password','secret')
    f = request.files.get('file')
    reader = PyPDF2.PdfReader(f)
    writer = PyPDF2.PdfWriter()
    for p in reader.pages:
        writer.add_page(p)
    writer.encrypt(user_pwd=pwd, owner_pwd=None)
    out = BytesIO()
    writer.write(out)
    out.seek(0)
    return send_file(out, as_attachment=True, download_name='protected.pdf', mimetype='application/pdf')

# Unlock PDF (if owner password known) — simple wrapper using pikepdf
@app.route('/unlock', methods=['POST'])
def unlock_pdf():
    pwd = request.form.get('password','')
    f = request.files.get('file')
    in_mem = BytesIO(f.read())
    out_mem = BytesIO()
    with pikepdf.open(in_mem, password=pwd) as pdf:
        pdf.save(out_mem)
    out_mem.seek(0)
    return send_file(out_mem, as_attachment=True, download_name='unlocked.pdf', mimetype='application/pdf')

# OCR (requires tesseract + pdf2image)
@app.route('/ocr', methods=['POST'])
def ocr_pdf():
    import pytesseract
    f = request.files.get('file')
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
    tmp.write(f.read())
    tmp.flush()
    pages = convert_from_path(tmp.name, dpi=200)
    text = '\n\n'.join(pytesseract.image_to_string(p) for p in pages)
    os.unlink(tmp.name)
    return jsonify({'text': text})

# Table extraction using camelot (if installed and file is vector PDF)
@app.route('/pdf2excel', methods=['POST'])
def pdf_to_excel():
    if camelot is None:
        return jsonify({'error':'camelot not installed on server'}), 500
    f = request.files.get('file')
    src = os.path.join(app.config['UPLOAD_FOLDER'], secure_filename(f.filename))
    f.save(src)
    tables = camelot.read_pdf(src, pages='all')
    out_xlsx = src + '.xlsx'
    # simple approach: write first table
    if len(tables) > 0:
        tables[0].to_excel(out_xlsx)
        return send_file(out_xlsx, as_attachment=True)
    else:
        return jsonify({'error':'no tables detected'}), 400

if __name__ == '__main__':
    app.run(debug=True)